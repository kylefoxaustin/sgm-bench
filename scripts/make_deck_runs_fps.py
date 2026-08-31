#!/usr/bin/env python3
"""A two-slide review deck: the Run 1 / Run 2 figures, frame-rate first."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
INK=RGBColor(0x20,0x21,0x24); MUTED=RGBColor(0x5F,0x63,0x68)
prs=Presentation(); prs.slide_width,prs.slide_height=Inches(13.333),Inches(7.5)
blank=prs.slide_layouts[6]
def tb(sl,x,y,w,h,t,size,bold=False,color=INK):
    b=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame; tf.word_wrap=True
    r=tf.paragraphs[0].add_run(); r.text=t
    r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name="Calibri"
SLIDES=[("Run 1 — 1920×800 input","docs/run1_fps.png",
         "Frame rate on a linear axis. Each bar carries its fps, and in parentheses the work rate (MDE/s) and the "
         "wall-clock milliseconds behind it."),
        ("Run 2 — 1920×1080 input","docs/run2_fps.png",
         "The same roster and the same algorithm at the other input size — the only variable between these two slides "
         "is the input size.")]
for title,png,note in SLIDES:
    s=prs.slides.add_slide(blank)
    # the figure carries its own title; a slide title would only repeat it
    h=6.45; w=h*(12.4/7.0)
    s.shapes.add_picture(png,Inches((13.333-w)/2),Inches(.22),height=Inches(h))
    tb(s,.6,6.82,12.2,.45,note,11.5,False,MUTED)
prs.save("docs/sgm-runs-fps.pptx")
print("wrote docs/sgm-runs-fps.pptx")
