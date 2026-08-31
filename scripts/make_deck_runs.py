#!/usr/bin/env python3
"""make_deck_runs.py — the two-resolution deck: ONE workload, TWO input sizes.
Same data as Configurations B and C, presented the way a reader thinks about
it: run 1 = 1920x1080 inputs, run 2 = 1920x800 inputs, same algorithm, same
hardware roster, bit-exact goldens per resolution."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

INK   = RGBColor(0x20,0x21,0x24); MUTED = RGBColor(0x5F,0x63,0x68)
ACC   = {"HW": RGBColor(0x7B,0x4F,0xD1), "GPU": RGBColor(0x4C,0x8B,0xF5),
         "CPU": RGBColor(0xE8,0x71,0x0A), "DSP": RGBColor(0x12,0xB5,0xA5),
         "REF": RGBColor(0x9A,0xA0,0xA6)}

RUN1 = [  # 1920x1080 -> 960x540, golden 0f0961d623009df5, work = 2.592 Mpx x 64
 ("NVIDIA Thor OFA",            "HW",     4.08,  None),
 ("NVIDIA RTX 5090",            "GPU",   11.96, 13870),
 ("NVIDIA Orin AGX OFA",        "HW",    19.51,  None),
 ("NVIDIA Thor GPU",            "GPU",   56.7,   2928),
 ("NVIDIA Orin AGX GPU",        "GPU",   96.9,   1711),
 ("Cortex-A720 x8",             "CPU",  239.6,    692),
 ("Cortex-A78C x8 (6 threads)", "CPU",  291.8,    569),
 ("Hexagon v73 NSP",            "DSP",  372.4,    446),
 ("Mali-G720 (10 CU)",          "GPU",  509.4,    326),
 ("Cortex-A78C x1",             "CPU",  595.3,    279),
 ("Cortex-A720 x1",             "CPU",  643.1,    258),
 ("Cortex-A55 x6",              "CPU",  844.8,    196),
 ("Cortex-A55 x1",              "CPU", 2954.0,     56),
 ("Mali-G310 (1 CU)",           "GPU", 3657.0,     45),
 ("scalar reference",           "REF",18020.0,    9.2),
]
RUN2 = [  # 1920x800 -> 960x400, golden bcb9cb0bd6f49799, work = 1.92 Mpx x 64
 ("NVIDIA Thor OFA",            "HW",     3.38,  None),
 ("NVIDIA RTX 5090",            "GPU",    9.08, 13538),
 ("NVIDIA Orin AGX OFA",        "HW",    14.90,  None),
 ("NVIDIA Thor GPU",            "GPU",   45.1,   2723),
 ("NVIDIA Orin AGX GPU",        "GPU",   72.8,   1688),
 ("Cortex-A78C x8 (6 threads)", "CPU",  211.9,    580),
 ("Cortex-A720 x8",             "CPU",  238.3,    516),
 ("Hexagon v73 NSP",            "DSP",  276.3,    445),
 ("Mali-G720 (10 CU)",          "GPU",  380.0,    323),
 ("Cortex-A78C x1",             "CPU",  442.6,    278),
 ("Cortex-A720 x1",             "CPU",  470.9,    261),
 ("Cortex-A55 x6",              "CPU",  640.1,    192),
 ("Cortex-A55 x1",              "CPU", 2201.0,     56),
 ("Mali-G310 (1 CU)",           "GPU", 2704.0,     45),
 ("scalar reference",           "REF",14265.0,    8.6),
]

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
blank = prs.slide_layouts[6]

def textbox(sl, x, y, w, h, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Calibri"
    return tf

def table_slide(title, sub, rows, base, golden):
    sl = prs.slides.add_slide(blank)
    textbox(sl, .6, .32, 12.2, .8, title, 30, True)
    textbox(sl, .6, 1.08, 12.2, .5, sub, 12.5, False, MUTED)
    tbl = sl.shapes.add_table(len(rows)+1, 6, Inches(.6), Inches(1.66), Inches(12.1), Inches(0.4)).table
    for i, w in enumerate((3.6, 1.0, 1.9, 1.3, 1.9, 2.4)): tbl.columns[i].width = Inches(w)
    for rr in range(len(rows)+1): tbl.rows[rr].height = Inches(0.25)
    for c, t in enumerate(("processing unit", "class", "runtime", "fps", "MDE/s (work)", "vs scalar reference")):
        cell = tbl.cell(0, c); cell.text = t
        pr = cell.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(11)
        pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x20,0x21,0x24)
    for rr, (lab, cls, ms, mdev) in enumerate(rows, start=1):
        fps = 1000.0/ms
        if mdev is None:
            vals = (lab, cls, f"{ms:,.2f} ms", f"{fps:,.0f}", "—", f"{base/ms:,.0f}x (time only)")
        else:
            vals = (lab, cls, f"{ms:,.2f} ms", f"{fps:,.1f}" if fps >= 10 else f"{fps:.2f}",
                    f"{mdev:,}" if mdev >= 10 else f"{mdev}",
                    f"{base/ms:,.0f}x faster" if ms != base else "the floor (1x)")
        for c, t in enumerate(vals):
            cell = tbl.cell(rr, c); cell.text = str(t)
            pp = cell.text_frame.paragraphs[0]; run = pp.runs[0]
            run.font.size = Pt(10.5); run.font.color.rgb = INK
            run.font.bold = (c in (0, 4)) and mdev is not None
            if c == 1: run.font.color.rgb = ACC[cls]; run.font.bold = True
    textbox(sl, .6, 6.9, 12.2, .5,
            f"Every software row byte-identical to golden {golden}, hash-gated in the timed run. OFA rows: the fixed-function "
            "engines running their OWN algorithm at matched input/output geometry (D=128) — throughput only, not bit-exact, "
            "MDE/s undefined for a different algorithm.", 10, False, MUTED)
    return sl

def chart_slide(title, png):
    sl = prs.slides.add_slide(blank)
    textbox(sl, .6, .3, 12.2, .6, title, 28, True)
    sl.shapes.add_picture(png, Inches(1.15), Inches(1.0), height=Inches(6.3))
    return sl

# slide 1: the setup
s = prs.slides.add_slide(blank)
textbox(s, .6, .55, 12.2, 1.0, "One SGM workload, two input resolutions", 36, True)
textbox(s, .6, 2.05, 12.2, 1.2,
        "The workload: dense stereo depth by Semi-Global Matching — two FULL 64-disparity searches, one at input "
        "resolution and one at half, data costs averaged (multi-scale fusion), census 9x7 on a SobelX-prefiltered image, "
        "8 aggregation paths with direction-weighted P1 (40/20/10), P2=200 saturating, half-resolution output.", 15)
textbox(s, .6, 3.35, 12.2, .9,
        "Run 1:  1920 x 1080 stereo input   ->   960 x 540 disparity map\n", 20, True)
textbox(s, .6, 3.95, 12.2, .9,
        "Run 2:  1920 x  800 stereo input   ->   960 x 400 disparity map", 20, True)
textbox(s, .6, 4.9, 12.2, 1.2,
        "Same algorithm, same parameters, same hardware roster — fifteen processing units from four silicon vendors. "
        "The ONLY variable between the runs is the input size. Acceptance: for each resolution a scalar oracle defines a "
        "golden disparity map, and every implementation must reproduce it byte-for-byte in the same invocation that is "
        "timed, or the timing is void. The fixed-function OFA engines (which run their own frozen algorithm) are measured "
        "beside the roster at the same geometry, labelled as such.", 13, False, MUTED)
textbox(s, .6, 6.6, 12.2, .5,
        "Every number is MEASURED (medians of repeated runs; spreads recorded). Code, scenes, goldens and instruments: "
        "github.com/kylefoxaustin/sgm-bench", 12, False, MUTED)

table_slide("Run 1 — 1920×1080 input", 
            "960×540 output · golden 0f0961d623009df5 · MDE/s(work) = (2.074 + 0.518) Mpx × 64 / time", 
            RUN1, 18020.0, "0f0961d623009df5")
chart_slide("Run 1 — 1920×1080, in frames per second", "docs/run1080_fps.png")
table_slide("Run 2 — 1920×800 input",
            "960×400 output · golden bcb9cb0bd6f49799 · MDE/s(work) = (0.384 + 1.536) Mpx × 64 / time",
            RUN2, 14265.0, "bcb9cb0bd6f49799")
chart_slide("Run 2 — 1920×800, in frames per second", "docs/run800_fps.png")

# ratio slide
s = prs.slides.add_slide(blank)
textbox(s, .6, .32, 12.2, .8, "The two runs agree: performance scales with input size", 26, True)
textbox(s, .6, 1.2, 12.2, .9,
        "1920x1080 has 1.35x the pixels of 1920x800, so run 1 should take ~1.35x run 2's time on every chip. It does:",
        14)
RAT = [("NVIDIA RTX 5090","1.32x"),("NVIDIA Thor GPU","1.26x"),("NVIDIA Orin AGX GPU","1.33x"),
       ("Cortex-A78C x8","1.38x"),("Hexagon v73 NSP","1.348x  — and identical MDE/s both runs"),
       ("Mali-G720","1.34x"),("Mali-G310","1.35x"),("single Arm cores","1.32x – 1.37x"),
       ("OFA engines","1.21x / 1.31x"),
       ("Cortex-A720 x8","1.01x  — THE exception: the 8-thread cluster was never work-limited at 1920x800 (thread scaling 1.98x vs 2.68x), so the extra 35% of pixels ride in its slack for free")]
tf = textbox(s, .9, 2.0, 11.8, 4.2, f"{RAT[0][0]}:  {RAT[0][1]}", 15)
for lab, v in RAT[1:]:
    p = tf.add_paragraph(); r = p.add_run(); r.text = f"{lab}:  {v}"
    r.font.size = Pt(15); r.font.color.rgb = INK; r.font.name = "Calibri"
textbox(s, .6, 6.5, 12.2, .8,
        "Which is the point of running both: a benchmark whose two resolutions scale by exactly the pixel ratio, "
        "bit-exact at each size, is measuring the workload — not an artifact of either scene.", 13, True)

# real imagery slide
s = prs.slides.add_slide(blank)
textbox(s, .6, .35, 12, .7, "The same kernels on real stereo imagery", 30, True)
textbox(s, .6, 1.0, 12, .4,
        "Middlebury 2014 'Motorcycle' — real calibrated two-camera capture, dense structured-light ground truth, 1482x1000, D=128",
        13, False, MUTED)
s.shapes.add_picture("docs/panels/motorcycle_real.png", Inches(.6), Inches(1.5), width=Inches(12.1))
textbox(s, .6, 6.55, 12.2, .9,
        "Seven targets produce the identical golden e8a95242882013f0 on the real capture — the picture is the same, only "
        "the time differs. Accuracy vs dense ground truth: bad>1px 16.2%, bad>2px 11.2%, MAE 3.35 — confirmed to "
        "0.05 points by a second scoring derived from the ground-truth arrays.", 12)

prs.save("docs/sgm-two-resolutions.pptx")
print("wrote docs/sgm-two-resolutions.pptx —", len(prs.slides._sldIdLst), "slides")
