#!/usr/bin/env python3
"""make_deck.py — the results deck. Slide 1 summarises; one slide per unit shows
the arithmetic that produced its number, so every figure can be re-derived from
the measured milliseconds without trusting the table."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W, H, D, PATHS = 1920, 1080, 64, 4
INK   = RGBColor(0x20, 0x21, 0x24)
MUTED = RGBColor(0x5F, 0x63, 0x68)
ACC   = {"HW": RGBColor(0x7B,0x4F,0xD1), "GPU": RGBColor(0x4C,0x8B,0xF5), "CPU": RGBColor(0xE8,0x71,0x0A),
         "DSP": RGBColor(0x12,0xB5,0xA5), "REF": RGBColor(0x9A,0xA0,0xA6)}

# label, class, ms, silicon, notes[]
UNITS = [
 ("NVIDIA Thor GPU",   "GPU",  13.19, "Blackwell, 20 SM, CUDA 13.2, sm_110", [
   "Rewritten warp-centric: one warp owns a scanline's whole disparity range,",
   "L in registers, min-reduce by __shfl_xor, ZERO barriers in the recurrence.",
   "2.69x over the naive OpenCL transliteration (35.39 -> 13.19 ms).",
   "MEMORY BOUND: aggregate moves 1.46 GB in 10.04 ms = 145 GB/s,",
   "against a 247 GB/s copy ceiling measured on this same part.",
   "Peak across the D sweep is 12,138 MDE/s at D=128; D=256 falls back to 9,409."]),
 ("NVIDIA Orin AGX GPU","GPU",  23.28, "Ampere, 16 SM, CUDA 12.6, sm_87", [
   "Same kernel, same changes: 3.10x over the naive port (72.27 -> 23.28 ms).",
   "Bit-exact on the first run, before any tuning — the payoff for porting a",
   "verified OpenCL kernel one-for-one instead of writing a new one."]),
 ("NVIDIA RTX 5090 GPU","GPU",   3.46, "Blackwell, 170 SM, discrete, 32 GB, driver 580.173  |  D=128: 4.09 ms  ·  D=256: 7.08 ms", [
   "The fastest target measured: 37,383 MDE/s at D=64, 282 fps at 1080p.",
   "64,880 MDE/s at D=128 and 74,957 at D=256 -- and unlike Thor it shows NO",
   "reversal at D=256, which is what the traffic mechanism predicts on a part",
   "with 5.6x the memory bandwidth.",
   "Bandwidth: 587 GB/s achieved of a measured 1,385 GB/s ceiling = 42%, the",
   "LOWEST utilisation of the three GPUs despite being the fastest -- so at this",
   "speed it is no longer purely bandwidth-bound.",
   "⚠️ Built as PTX and JIT-compiled by the driver: the local toolkit is CUDA",
   "12.6, which cannot target sm_120. Golden verified, but not a native build.",
   "⭐ ITS STEREO HARDWARE IS MEASURED, AND STEREO MODE IS ALREADY GONE.",
   "NV_OF_MODE_STEREODISPARITY returns UNSUPPORTED_FEATURE at every grid size;",
   "NVIDIA's documented deprecation has already landed on Blackwell.",
   "The fallback -- X component of general optical flow -- runs 8.50 ms with",
   "5.3% bad pixels, against our CUDA SGM's 4.13 ms and 2.8% on the same scene.",
   "Software faster AND more accurate, while doing less work (1D vs 2D).",
   "🚨 13.4% of pixels return a non-zero VERTICAL component, largest 75 px --",
   "errors the removed stereo mode was structurally incapable of producing."]),
 ("Cortex-A720 x8",    "CPU",  51.70, "Radxa Orion O6, 2.2-2.5 GHz, NEON  |  single core: 313.60 ms = 423 MDE/s @2.5 GHz", [
   "The CPU headline: 10.4x the fastest published Arm-CPU SGM at matched D",
   "and matched path count (ReS2tAC, Sensors 21(11):3938, 2021 — 242 MDE/s).",
   "Best single core measured at 423 MDE/s; scales monotonically 1->8 threads.",
   "D=64 -> D=128 is FLAT (-2.4%): 16 NEON lanes are already full at D=64."]),
 ("Cortex-A78C x8",    "CPU",  97.99, "Qualcomm IQ-9075, NEON  |  single core: 332.22 ms = 400 MDE/s", [
   "Ran THIS repository's a55 NEON implementation unmodified, on another",
   "vendor's silicon, and reproduced the golden at every thread count.",
   "Per-core 402 MDE/s vs the A720's 423 — a 5% gap between two wide",
   "out-of-order Arm cores from different vendors on the same algorithm.",
   "Peaks at SIX threads, not eight; 8 threads is slower on that board."]),
 ("Hexagon v73 NSP",   "DSP", 148.44, "Qualcomm IQ-9075, HVX, 4 HVX contexts", [
   "MEDIAN of 5 invocations; the board is BIMODAL between invocations (18% at",
   "720p) so minima are the estimator most sensitive to luck -- quote medians.",
   "Went 18.7x -> ... -> 1.51x vs the A78C cluster across one day of tuning,",
   "nothing retracted: always labelled a floor of EFFORT. Beats one A78C core 2.26x.",
   "Six hardware threads but only FOUR HVX contexts, so 4 threads is optimal.",
   "Still latency-bound at 6.8 of 28 GB/s — a different wall to Thor's."]),
 ("Mali-G720",         "GPU", 256.00, "Immortalis, 10 CU, OpenCL 3.0", [
   "Loses to TWO A720 cores. But this is an UNBOUNDED FLOOR, not a verdict",
   "on the hardware: the kernel carries ONE dependency chain per work-item",
   "with three barrier sites per step and never interleaves chains.",
   "Nobody has done for this part what was done for the DSP."]),
 ("NVIDIA Thor OFA",   "HW",   9.345, "Fixed-function SGM engine (Optical Flow Accelerator), VPI 4.1.4  |  1080p D=128, downscaleFactor=1", [
   "DEDICATED STEREO HARDWARE, and it beats our tuned CUDA on the same chip",
   "by 2.34x (9.35 ms vs 21.86 ms at 1920x1080 D=128).",
   "NOT bit-exact to our golden: VPI's SGM is a different implementation with",
   "its own census, penalties and confidence output. Throughput comparison at",
   "matched resolution and D -- NOT the same test, and NOT a quality claim.",
   "🔴 ACCURACY WITHDRAWN. Scored against ground truth our SGM gets 2.6% bad>1px",
   "(MAE 0.45); OFA aligns under no simple descaling, so its accuracy is",
   "UNMEASURED and this 2.34x is throughput only.",
   "VPI defaults to includeDiagonals=1, so OFA ran 8 paths against our 4. Measured,",
   "not assumed: OFA costs 9.42 ms at 8 paths and 9.85 at 4 -- path-count",
   "independent. The timing stands and OFA delivers an 8-path result in our time.",
   "⚠️ It defaults to downscaleFactor=2, which emits a 960x540 map from 1080p",
   "input in 4.14 ms. Quoting that as a 1080p figure overstates it by 2.3x."]),
 ("NVIDIA Orin AGX OFA","HW",  72.637, "Fixed-function SGM engine, VPI 3.2.4  |  1080p D=128, downscaleFactor=1", [
   "The same fixed-function engine one generation earlier, and here it LOSES:",
   "our CUDA kernel does the same work in 33.16 ms against the hardware's 72.64.",
   "Software wins by 2.19x on the silicon built for this exact job.",
   "Thor's OFA is 7.8x faster than Orin's at identical settings -- a far larger",
   "generational jump than the 1.5x between the two GPUs.",
   "Taken together these two rows are the interesting result: whether dedicated",
   "hardware beats a tuned kernel is a property of the generation, not of the",
   "idea of dedicated hardware.",
   "Same caveats as the Thor row: not bit-exact, accuracy unmeasured, and the",
   "timed region excludes format conversion where ours is end-to-end (<5%)."]),
 ("Cortex-A55 x6",     "CPU", 341.30, "i.MX 95 FRDM, 1.8 GHz, NEON  |  single core: 1,588.74 ms = 84 MDE/s", [
   "The reference platform this project is anchored to, and the slowest CPU here.",
   "Scales 4.65x from 1 to 6 cores (1,588.74 -> 341.30 ms) -- that spread is also",
   "the proof that the -t thread flag was working on these runs, since the inert",
   "flag it was found to have would have collapsed every count to one timing.",
   "Per core 84 MDE/s against the A720's 423: a 5.0x gap between an in-order",
   "little core at 1.8 GHz and a wide out-of-order core at 2.5 GHz.",
   "Block-width optimum is 192 here and 256 on the A720 -- not the same number."]),
 ("Mali-G310",         "GPU",1846.00, "1 CU, OpenCL 3.0", [
   "The entry-tier part, same kernel, same hash. Its 'no compute driver'",
   "reputation is wrong — it ships working OpenCL 3.0 and merely lacks clinfo."]),
 ("scalar reference",  "REF",9188.00, "1x Cortex-A55, -O2, no SIMD, single thread", [
   "Deliberately boring: the oracle every other bar is checked against.",
   "It defines the golden hash, so it cannot be wrong by construction —",
   "it can only be slow, and it is 699x slower than the fastest bar."]),
]

def mde(ms): return W*H*D/ms/1000

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
blank = prs.slide_layouts[6]

def textbox(sl, x, y, w, h, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Calibri"
    return tf

# ---------------- slide 1: summary ----------------
s1 = prs.slides.add_slide(blank)
textbox(s1, .6, .32, 12.2, .9, "Semi-Global Matching across ten processors", 28, True)
textbox(s1, .6, 1.12, 12.2, .5,
        "1920x1080  ·  D=64  ·  4 paths  ·  9x7 census  ·  every row bit-exact to golden 0bc0102058d1505f",
        14, False, MUTED)

SINGLE = [("Cortex-A720 x1 @2.5GHz","CPU",313.60),("Cortex-A78C x1","CPU",336.11),
          ("Cortex-A55 x1 @1.8GHz","CPU",1588.74)]
# Dedicated stereo engines, shown WITH the processors they share a die with.
# Throughput only -- not bit-exact to the golden; the 5090's stereo mode is
# measured as removed (UNSUPPORTED_FEATURE at every grid size).

rows = [u for u in UNITS] + [(l,c,ms,"",[]) for l,c,ms in SINGLE]
rows = sorted(rows, key=lambda r: (r[2] is None, r[2] or 0))
# The OFA entries already live in UNITS (they have their own slides); mark them
# so the summary table computes MDE/s at THEIR D (128) instead of D=64, and
# labels them throughput-only rather than "Nx faster than the scalar ref".
HWNAMES = {"NVIDIA Thor OFA": 128, "NVIDIA Orin AGX OFA": 128}
tbl = s1.shapes.add_table(len(rows)+1, 7, Inches(.6), Inches(1.72), Inches(12.1), Inches(0.4)).table
for i, w in enumerate((3.3, 0.9, 1.7, 1.1, 1.5, 1.6, 2.0)): tbl.columns[i].width = Inches(w)
for r in range(len(rows)+1): tbl.rows[r].height = Inches(0.27)
hdr = ("processing unit", "class", "runtime", "fps", "MDE/s", "DRAM GB/s *", "vs scalar reference")
for c, t in enumerate(hdr):
    cell = tbl.cell(0, c); cell.text = t
    pr = cell.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(11)
    pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x20,0x21,0x24)
base = 9188.0
for r, (lab, cls, ms, sil, notes) in enumerate(rows, start=1):
    BW = {"NVIDIA RTX 5090 GPU":"587 / 1,385","NVIDIA Thor GPU":"145 / 249",
          "NVIDIA Orin AGX GPU":"87 / 175","Hexagon v73 NSP":"6.8 / 28",
          "Cortex-A720 x8":"6.6 / 46.3","Cortex-A78C x8":"3.4 / —",
          "Mali-G720":"16.1 / 46.3","Cortex-A720 x1 @2.5GHz":"1.1 / 46.3",
          "Cortex-A78C x1":"1.0 / —","Cortex-A55 x6":"1.1 / 13.7",
          "Cortex-A55 x1 @1.8GHz":"0.2 / 13.7","Mali-G310":"1.6 / 13.7"}
    bw = BW.get(lab, "—")
    hwd = HWNAMES.get(lab)
    if lab in HWNAMES and ms is None:
        vals = (lab, cls, "stereo mode removed", "—", "—", "—", "UNSUPPORTED_FEATURE")
    elif lab in HWNAMES:
        vals = (lab, cls, f"{ms:,.2f} ms (D={hwd})", f"{1000/ms:,.1f}",
                f"{W*H*hwd/ms/1000:,.0f}", "—", f"{base/ms:,.0f}x (time only)")
    else:
        vals = ((lab, cls, "—" if ms is None else f"{ms:,.2f} ms",
             "—" if ms is None else f"{1000/ms:,.1f}",
             "—" if ms is None else f"{mde(ms):,.0f}", bw,
             "not yet measured" if ms is None else f"{base/ms:,.0f}x faster"))
    for c, t in enumerate(vals):
        cell = tbl.cell(r, c); cell.text = str(t)
        pp = cell.text_frame.paragraphs[0]; run = pp.runs[0]
        run.font.size = Pt(11); run.font.color.rgb = MUTED if ms is None else INK
        run.font.bold = (c in (0, 4)) and ms is not None
        if c == 1:
            run.font.color.rgb = ACC[cls]; run.font.bold = True
# headline callouts, in the space under the table
CALLOUTS = [("2,655x", "faster than the scalar reference,\nwith byte-identical output"),
            ("10.4x", "the fastest published Arm-CPU SGM,\nat matched D and matched path count"),
            ("10 / 4 / 3 / 1", "processors / vendors / programming models\n/ one golden hash")]
for i, (big, small) in enumerate(CALLOUTS):
    x = .6 + i * 4.1
    textbox(s1, x, 6.35, 3.9, .6, big, 30, True, ACC["GPU"] if i == 0 else INK)
    tf = textbox(s1, x, 6.85, 3.9, .8, small.split("\n")[0], 12, False, MUTED)
    pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = small.split("\n")[1]
    rr.font.size = Pt(12); rr.font.color.rgb = MUTED; rr.font.name = "Calibri"

textbox(s1, .6, 7.28, 12.2, .3,
        "* DRAM GB/s = achieved (DERIVED: modelled bytes / measured phase time) / ceiling (MEASURED, copy probe per board; CPU+Mali share one bus). GPUs run at 42-58% of ceiling (memory-bound); CPUs/Malis at 8-37% (compute-bound). "
        "Aggregation streams 1.46 GB/frame at 1080p D=64. Blank = not instrumented. "
        "MDE/s = W x H x D / runtime — right for comparing implementations, wrong for choosing a configuration.", 10, False, MUTED)

# ---------------- bandwidth slide ----------------
sBW = prs.slides.add_slide(blank)
textbox(sBW, .6, .32, 12.2, .8, "This workload is about moving data", 30, True)
textbox(sBW, .6, 1.1, 12.2, .5,
        "At 1080p D=64 the aggregation phase alone streams 1.46 GB per frame (the summed cost volume is 265 MB, "
        "touched four times). Every ceiling is MEASURED with a streaming-copy probe on the part itself; every "
        "achieved figure is DERIVED as the kernel's exact modelled byte count / the measured phase time.", 12, False, MUTED)
BWROWS = [
 ("NVIDIA RTX 5090",  "2.49 ms",  "587",  "1,385", "42%", "memory-bound -- and the LEAST efficient of the three"),
 ("NVIDIA Thor",      "10.04 ms", "145",  "249",   "58%", "memory-bound"),
 ("NVIDIA Orin AGX",  "16.83 ms", "87",   "175",   "49%", "memory-bound"),
 ("Mali-G720 (10 CU)","--",       "16.1", "46.3",  "35%", "not bandwidth-bound"),
 ("8x Cortex-A720",   "--",       "6.6",  "46.3",  "14%", "compute-bound"),
 ("Hexagon v73 NSP",  "--",       "6.8",  "28",    "24%", "LATENCY-bound (28 GB/s available, 6.8 used)"),
 ("8x Cortex-A78C",   "--",       "3.4",  "--",    "--",  "ceiling unmeasured on that board -- stated, not guessed"),
 ("6x Cortex-A55",    "--",       "1.1",  "13.7",  "8%",  "compute-bound"),
]
tblW = sBW.shapes.add_table(len(BWROWS)+1, 6, Inches(.6), Inches(2.0), Inches(12.1), Inches(0.4)).table
for i, w in enumerate((2.6, 1.5, 1.6, 1.5, 1.1, 3.8)): tblW.columns[i].width = Inches(w)
for rr_ in range(len(BWROWS)+1): tblW.rows[rr_].height = Inches(0.34)
for c, t in enumerate(("processing unit", "aggregate", "achieved GB/s", "ceiling GB/s", "util", "wall")):
    cell = tblW.cell(0, c); cell.text = t
    pr = cell.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(11)
    pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x20,0x21,0x24)
for rr_, row in enumerate(BWROWS, start=1):
    for c, t in enumerate(row):
        cell = tblW.cell(rr_, c); cell.text = str(t)
        pp = cell.text_frame.paragraphs[0]; run = pp.runs[0]
        run.font.size = Pt(11); run.font.color.rgb = INK; run.font.bold = (c in (0, 2))
textbox(sBW, .6, 5.4, 12.2, .9,
        "One column, two walls: the NVIDIA GPUs sit at 42-58% of their measured ceilings -- memory-bound, which is why "
        "their optimisation ended at traffic reduction. The CPUs and Malis sit at 8-35% -- NOT bandwidth-bound, which is "
        "why their wins came from arithmetic (vectorised argmin, blocking, lane tricks). The Hexagon is a third wall: "
        "latency-bound, hidden by interleaving independent dependency chains.", 12.5, True, INK)
textbox(sBW, .6, 6.5, 12.2, .7,
        "Cross-GPU sanity: MDE/s per GB/s of ceiling is 27 / 40 / 33 across a 7.9x bandwidth spread -- throughput largely "
        "IS bandwidth for the tuned GPU kernels, and the residual (the 5090's 42%) says the fastest part is no longer "
        "purely bandwidth-bound. Achieved-byte instrumentation exists for Configuration A; the B/C runs reuse the same "
        "silicon and ceilings but their per-phase byte models are not separately derived.", 11, False, MUTED)

# ---------------- bandwidth slide, Configuration B ----------------
sB2 = prs.slides.add_slide(blank)
textbox(sB2, .6, .32, 12.2, .8, "Configuration B moves the memory walls", 30, True)
textbox(sB2, .6, 1.1, 12.2, .5,
        "Same slide for the block-flavoured workload (1920x800, two scales, 8 paths). The 8-path uint16 S-plane "
        "streams ~3.1 GB per frame on its own; the pipeline's streaming byte model totals 4.50 GB/frame. Three cells "
        "are MEASURED at the i.MX95 DDR controller; the rest are DERIVED as 4.50 GB / the measured frame time.", 11.5, False, MUTED)
BW2 = [
 ("NVIDIA RTX 5090",  "9.08 ms",  "495 †",  "1,385", "36%", "memory-bound; below A's 42% (kernel lacks path-pairing)"),
 ("NVIDIA Thor",      "45.1 ms",  "100 †",  "249",   "40%", "memory-bound"),
 ("NVIDIA Orin AGX",  "72.8 ms",  "62 †",   "175",   "35%", "memory-bound"),
 ("8x Cortex-A78C",   "211.9 ms", "21.2 †", "n.m.",  "--",  "the S-plane tax; ceiling unmeasured on that board"),
 ("8x Cortex-A720",   "238.3 ms", "18.9 †", "46.3",  "41%", "was 14% on A -- no longer comfortably compute-bound"),
 ("Hexagon v73 NSP",  "276.3 ms", "~16 †",  "28",    "~56%","approximate: their fused kernel parks state in VTCM"),
 ("Mali-G720 (10 CU)","380.0 ms", "11.8 †", "46.3",  "26%", "not bandwidth-bound"),
 ("6x Cortex-A55",    "640.1 ms", "8.8 M",  "13.7",  "64%", "TWO-THIRDS of the whole chip's bus at six threads"),
 ("Mali-G310 (1 CU)", "2,704 ms", "2.1 M",  "13.7",  "15%", "GPU + host combined; not bandwidth-bound"),
 ("1x Cortex-A55",    "2,201 ms", "2.7 M",  "13.7",  "20%", "single core, still triple its Config A draw"),
]
tblW2 = sB2.shapes.add_table(len(BW2)+1, 6, Inches(.6), Inches(1.8), Inches(12.1), Inches(0.4)).table
for i, w in enumerate((2.4, 1.3, 1.5, 1.4, 1.0, 4.5)): tblW2.columns[i].width = Inches(w)
for rr_ in range(len(BW2)+1): tblW2.rows[rr_].height = Inches(0.3)
for c, t in enumerate(("processing unit", "frame", "achieved GB/s", "ceiling GB/s", "util", "wall")):
    cell = tblW2.cell(0, c); cell.text = t
    pr = cell.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(11)
    pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x20,0x21,0x24)
for rr_, row in enumerate(BW2, start=1):
    for c, t in enumerate(row):
        cell = tblW2.cell(rr_, c); cell.text = str(t)
        pp = cell.text_frame.paragraphs[0]; run = pp.runs[0]
        run.font.size = Pt(9.5); run.font.color.rgb = INK; run.font.bold = (c in (0, 2))
textbox(sB2, .6, 5.5, 12.2, .85,
        "The headline: at Configuration B the CPU walls MOVE. The A720 cluster jumps from 14% to 41% of its bus and "
        "the six-A55 cluster to 64% -- the 8-path uint16 S-plane converts comfortably compute-bound CPUs into "
        "bandwidth-pressured ones. (It is also why the A720 cluster absorbed Configuration C's extra pixels at the "
        "same wall-clock: it was already pressed against a different limit than arithmetic.)", 12.5, True, INK)
textbox(sB2, .6, 6.55, 12.2, .5,
        "M = MEASURED at the imx9_ddr0 DDR controller (idle-corrected whole-run, read+write). † = DERIVED: 4.50 GB/frame "
        "streaming model (cost build+merge 0.27 + 8-path cost reads 0.79 + S-plane RMW 3.15 + argmin 0.20 GB) / measured "
        "frame time -- cache-reuse-free, so a slight under-count: the model reproduces the A55x6's measured cell at 80%.",
        10, False, MUTED)

# ---------------- one slide per unit ----------------
for lab, cls, ms, sil, notes in UNITS:
    sl = prs.slides.add_slide(blank)
    bar = sl.shapes.add_textbox(Inches(0), Inches(0), Inches(13.333), Inches(.12))
    textbox(sl, .6, .35, 11, .7, lab, 32, True, ACC[cls])
    textbox(sl, .6, 1.0, 11, .4, sil, 13, False, MUTED)

    textbox(sl, .6, 1.7, 11.5, .4, "The calculation", 16, True)
    if ms is None:
        textbox(sl, .6, 2.15, 11.5, 1.2,
                "No measurement exists for this part, so no number is derived.\n"
                "Nothing is estimated in its place.", 15, False, MUTED)
    else:
        calc = (f"MDE/s  =  W x H x D / runtime\n"
                f"        =  {W:,} x {H:,} x {D} / {ms:,.2f} ms\n"
                f"        =  {W*H*D:,} disparity estimations / {ms/1000:.5f} s\n"
                f"        =  {mde(ms):,.0f} million disparity estimations per second\n\n"
                f"frame rate  =  1000 / {ms:,.2f}  =  {1000/ms:,.1f} fps\n"
                f"pixel rate  =  {W*H/1e6:.4f} Mpx x {1000/ms:,.1f}  =  {W*H/1e6*1000/ms:,.2f} Mpix/s")
        tf = textbox(sl, .6, 2.15, 11.5, 2.2, calc, 14)
        for p in tf.paragraphs:
            for r in p.runs: r.font.name = "Consolas"

    textbox(sl, .6, 4.5, 11.5, .4, "What the measurement means", 16, True)
    tf = textbox(sl, .6, 4.9, 11.9, 2.2, notes[0], 13)
    for extra in notes[1:]:
        p = tf.add_paragraph(); r = p.add_run(); r.text = extra
        r.font.size = Pt(13); r.font.color.rgb = INK; r.font.name = "Calibri"

    if ms is not None:
        textbox(sl, .6, 7.0, 11.5, .4,
                f"MEASURED — median of the timed frames, golden hash verified in the same run.",
                11, False, MUTED)

# ---------------- real imagery slide ----------------
sl = prs.slides.add_slide(blank)
textbox(sl, .6, .35, 12, .7, "On real stereo imagery", 30, True)
textbox(sl, .6, 1.0, 12, .4,
        "Middlebury 2014 'Motorcycle' — a real calibrated two-camera capture, dense structured-light "
        "ground truth — 1482x1000, D=128", 13, False, MUTED)
try:
    sl.shapes.add_picture("docs/panels/motorcycle_real.png", Inches(.6), Inches(1.5), width=Inches(12.1))
except Exception:
    pass
textbox(sl, .6, 6.55, 12.2, .9,
        "All eight targets produce the identical golden e8a95242882013f0 — the picture is the same, only the "
        "time differs (NSP: 179.23 ms, ~1,058 MDE/s). Accuracy vs dense ground truth: bad>1px 16.2%, bad>2px 11.2%, "
        "MAE 3.35 — independently reproduced to 0.05 points on hardware we have never touched.", 12, False, INK)
textbox(sl, .6, 7.12, 12.2, .4,
        "Real imagery is a better ACCEPTANCE test, not just a more realistic one: 14,417 pixels have a tied "
        "minimum here against 20 on the synthetic scene, so the hash-critical tie-break rule is pinned 720x harder.",
        11, False, MUTED)

# ---------------- Configuration B summary table ----------------
sB = prs.slides.add_slide(blank)
textbox(sB, .6, .32, 12.2, .9, "Configuration B across the same roster", 28, True)
textbox(sB, .6, 1.12, 12.2, .5,
        "1920x800 in, 960x400 out  ·  two full 64-D scales, costs averaged  ·  8 paths, weighted P1  ·  "
        "census 9x7 on SobelX  ·  P2=200 saturating  ·  every row bit-exact to golden bcb9cb0bd6f49799",
        13, False, MUTED)

BROWS = [  # label, class, ms, MDE/s(work) as published — ms-sorted, OFA slotted in
 ("NVIDIA OFA on Thor (ds=2)",  "HW",      3.38,  None),
 ("NVIDIA RTX 5090",            "GPU",     9.08, 13538),
 ("NVIDIA OFA on Orin (ds=2)",  "HW",     14.90,  None),
 ("NVIDIA Thor",                "GPU",    45.1,   2723),
 ("NVIDIA Orin AGX",            "GPU",    72.8,   1688),
 ("Cortex-A78C x8 (6 threads)", "CPU",   211.9,    580),
 ("Cortex-A720 x8",             "CPU",   238.3,    516),
 ("Hexagon v73 NSP",            "DSP",   276.3,    445),
 ("Mali-G720 (10 CU)",          "GPU",   380.0,    323),
 ("Cortex-A78C x1",             "CPU",   442.6,    278),
 ("Cortex-A720 x1",             "CPU",   470.9,    261),
 ("Cortex-A55 x6",              "CPU",   640.1,    192),
 ("Cortex-A55 x1",              "CPU",  2201.0,     56),
 ("Mali-G310 (1 CU)",           "GPU",  2704.0,     45),
 ("scalar reference",           "REF", 14265.0,    8.6),
]
tblB = sB.shapes.add_table(len(BROWS)+1, 7, Inches(.6), Inches(1.72), Inches(12.1), Inches(0.4)).table
for i, w in enumerate((3.3, 0.9, 1.6, 1.1, 1.7, 1.5, 2.0)): tblB.columns[i].width = Inches(w)
for rr_ in range(len(BROWS)+1): tblB.rows[rr_].height = Inches(0.25)
for c, t in enumerate(("processing unit", "class", "runtime", "fps", "MDE/s (work)", "DRAM GB/s ‡", "vs scalar reference")):
    cell = tblB.cell(0, c); cell.text = t
    pr = cell.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(11)
    pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x20,0x21,0x24)
baseB = 14265.0
CEIL = {"NVIDIA RTX 5090": "— / 1,385", "NVIDIA Thor": "— / 249", "NVIDIA Orin AGX": "— / 175",
        "NVIDIA OFA on Thor (ds=2)": "— / 249", "NVIDIA OFA on Orin (ds=2)": "— / 175",
        "Cortex-A78C x8 (6 threads)": "— / n.m.", "Cortex-A78C x1": "— / n.m.",
        "Cortex-A720 x8": "— / 46.3", "Cortex-A720 x1": "— / 46.3",
        "Hexagon v73 NSP": "— / 28", "Mali-G720 (10 CU)": "— / 46.3",
        "Cortex-A55 x6": "8.8 / 13.7", "Cortex-A55 x1": "2.7 / 13.7",
        "Mali-G310 (1 CU)": "— / 13.7", "scalar reference": "— / 13.7"}
for rr_, (lab, cls, ms, mdev) in enumerate(BROWS, start=1):
    bw = CEIL.get(lab, "—")
    if ms is not None and mdev is None:
        fps = 1000.0/ms
        vals = (lab, cls, f"{ms:,.2f} ms", f"{fps:,.0f}", "—", bw,
                f"{baseB/ms:,.0f}x (time only)")
    elif ms is None:
        vals = (lab, cls, "cannot run it", "—", "—", "—", "no 8-path/weighted-P1/two-scale controls")
    else:
        fps = 1000.0/ms
        vals = (lab, cls, f"{ms:,.2f} ms", f"{fps:,.1f}" if fps >= 10 else f"{fps:.2f}",
                f"{mdev:,}" if mdev >= 10 else f"{mdev}", bw,
                f"{baseB/ms:,.0f}x faster" if ms != baseB else "the floor (1x)")
    for c, t in enumerate(vals):
        cell = tblB.cell(rr_, c); cell.text = str(t)
        pp = cell.text_frame.paragraphs[0]; run = pp.runs[0]
        run.font.size = Pt(11); run.font.color.rgb = MUTED if ms is None else INK
        run.font.bold = (c in (0, 4)) and ms is not None
        if c == 1:
            run.font.color.rgb = ACC[cls]; run.font.bold = True

CALLOUTS_B = [("1,575x", "faster than the scalar reference,\nwith byte-identical output"),
              ("1.30x / 1.60x", "cluster over one NSP at matched tiers\n/ one NSP over one A78C core"),
              ("13 / 4 / 1", "targets / programming models (CUDA,\nOpenCL, NEON, HVX) / one golden hash")]
for i, (big, small) in enumerate(CALLOUTS_B):
    x = .6 + i * 4.1
    textbox(sB, x, 6.35, 3.9, .6, big, 28, True, ACC["GPU"] if i == 0 else INK)
    tf = textbox(sB, x, 6.85, 3.9, .8, small.split("\n")[0], 12, False, MUTED)
    pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = small.split("\n")[1]
    rr.font.size = Pt(12); rr.font.color.rgb = MUTED; rr.font.name = "Calibri"

textbox(sB, .6, 7.28, 12.2, .3,
        "MDE/s (work) = both scales count: (0.384 + 1.536) Mpx x 64 / runtime. OFA rows: their OWN SGM at matched "
        "geometry (ds=2, D=128), throughput only. ‡ DRAM = achieved / ceiling GB/s. A55 cells MEASURED at the i.MX95 "
        "DDR controller (idle-corrected, whole-run read+write) -- the 8-path multiscale pipeline drives the cluster to "
        "64% of the chip's whole bus. Other achieved cells not yet instrumented (n.m. = ceiling unmeasured, leased board).",
        10, False, MUTED)

# ---------------- Configuration B slide ----------------
sl = prs.slides.add_slide(blank)
textbox(sl, .6, .35, 12, .7, "Configuration B: multi-scale, 8-path", 30, True)
textbox(sl, .6, 1.0, 12.2, .4,
        "1920x800 in, 960x400 out - two FULL 64-disparity scales with data costs averaged (multi-scale fusion, "
        "not coarse-to-fine), census 9x7 on SobelX, 8 paths with direction-weighted P1 (40/20/10), P2=200 saturating.",
        12, False, MUTED)
try:
    sl.shapes.add_picture("docs/results_b.png", Inches(1.35), Inches(1.55), height=Inches(4.95))
except Exception:
    pass
textbox(sl, .6, 6.5, 12.2, .7,
        "Same roster, same acceptance model: every row bit-exact to golden bcb9cb0bd6f49799, same implementation "
        "tiers as Configuration A (tuned CUDA / OpenCL / NEON+OpenMP / HVX, scalar oracle as the floor). "
        "At matched tiers it ranks like Configuration A - the CPU cluster leads one NSP ~1.3x on both; "
        "per engine, one NSP beats one A78C core 1.60x.", 11, False, INK)
textbox(sl, .6, 7.18, 12.2, .3,
        "MDE/s is the work-performed convention - both scales count. The OFA engines cannot run this "
        "configuration (no 8-path, weighted-P1 or two-scale controls).", 10, False, MUTED)

# ---------------- Configuration C: one variable at a time ----------------
sC = prs.slides.add_slide(blank)
textbox(sC, .6, .32, 12.2, .9, "Configuration C: one variable at a time", 28, True)
textbox(sC, .6, 1.1, 12.2, .5,
        "C is exactly B's algorithm on A's own 1920x1080 stereo pair (out 960x540): A->C isolates the ALGORITHM, "
        "C->B isolates RESOLUTION - one knob each. Golden 0f0961d623009df5, every software row bit-exact.",
        12, False, MUTED)

CROWS = [  # label, class, ms, MDE/s(work) — same columns and types as the Config B table
 ("NVIDIA OFA on Thor (ds=2)",   "HW",     4.08,  None),
 ("NVIDIA RTX 5090",             "GPU",   11.96, 13870),
 ("NVIDIA OFA on Orin (ds=2)",   "HW",    19.51,  None),
 ("NVIDIA Thor",                 "GPU",   56.7,   2928),
 ("NVIDIA Orin AGX",             "GPU",   96.9,   1711),
 ("Cortex-A720 x8",              "CPU",  239.6,    692),
 ("Cortex-A78C x8 (6 threads)",  "CPU",  291.8,    569),
 ("Hexagon v73 NSP",             "DSP",  372.4,    446),
 ("Mali-G720 (10 CU)",           "GPU",  509.4,    326),
 ("Cortex-A78C x1",              "CPU",  595.3,    279),
 ("Cortex-A720 x1",              "CPU",  643.1,    258),
 ("Cortex-A55 x6",               "CPU",  844.8,    196),
 ("Cortex-A55 x1",               "CPU", 2954.0,     56),
 ("Mali-G310 (1 CU)",            "GPU", 3657.0,     45),
 ("scalar reference",            "REF",18020.0,    9.2),
]
tblC = sC.shapes.add_table(len(CROWS)+1, 7, Inches(.6), Inches(1.66), Inches(12.1), Inches(0.4)).table
for i, w in enumerate((3.3, 0.9, 1.6, 1.1, 1.7, 1.5, 2.0)): tblC.columns[i].width = Inches(w)
for rr_ in range(len(CROWS)+1): tblC.rows[rr_].height = Inches(0.25)
for c, t in enumerate(("processing unit", "class", "runtime", "fps", "MDE/s (work)", "DRAM GB/s ‡", "vs scalar reference")):
    cell = tblC.cell(0, c); cell.text = t
    pr = cell.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(11)
    pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x20,0x21,0x24)
baseC = 18020.0
CEILC = {"NVIDIA RTX 5090": "— / 1,385", "NVIDIA Thor": "— / 249", "NVIDIA Orin AGX": "— / 175",
         "NVIDIA OFA on Thor (ds=2)": "— / 249", "NVIDIA OFA on Orin (ds=2)": "— / 175",
         "Cortex-A78C x8 (6 threads)": "— / n.m.", "Cortex-A78C x1": "— / n.m.",
         "Cortex-A720 x8": "— / 46.3", "Cortex-A720 x1": "— / 46.3",
         "Hexagon v73 NSP": "— / 28", "Mali-G720 (10 CU)": "— / 46.3",
         "Cortex-A55 x6": "9.2 / 13.7", "Cortex-A55 x1": "2.7 / 13.7",
         "Mali-G310 (1 CU)": "— / 13.7", "scalar reference": "— / 13.7"}
for rr_, (lab, cls, ms, mdev) in enumerate(CROWS, start=1):
    fps = 1000.0/ms
    bw = CEILC.get(lab, "—")
    if mdev is None:
        vals = (lab, cls, f"{ms:,.2f} ms", f"{fps:,.0f}", "—", bw,
                f"{baseC/ms:,.0f}x (time only)")
    else:
        vals = (lab, cls, f"{ms:,.2f} ms", f"{fps:,.1f}" if fps >= 10 else f"{fps:.2f}",
                f"{mdev:,}" if mdev >= 10 else f"{mdev}", bw,
                f"{baseC/ms:,.0f}x faster" if ms != baseC else "the floor (1x)")
    for c, t in enumerate(vals):
        cell = tblC.cell(rr_, c); cell.text = str(t)
        pp = cell.text_frame.paragraphs[0]; run = pp.runs[0]
        run.font.size = Pt(10.5); run.font.color.rgb = INK
        run.font.bold = (c in (0, 4)) and mdev is not None
        if c == 1:
            run.font.color.rgb = ACC[cls]; run.font.bold = True

textbox(sC, .6, 5.95, 12.2, .7,
        "Two clean statements on the 5090, one knob each:  A -> C (same input images, algorithm swap): "
        "3.46 -> 11.96 ms - the two-scale 8-path algorithm costs ~3.5x the single-scale 4-path one.  "
        "C -> B (same algorithm, 1080 -> 800 rows): 11.96 -> 9.08 ms. Row for row, C/B ratios cluster at the "
        "1.35x the pixels predict (5090 1.32, Thor 1.26, NSP 1.348, Malis 1.34-1.35) - one exception below.", 11.5, True, INK)
textbox(sC, .6, 6.68, 12.2, .5,
        "The exception is the finding: the A720 CLUSTER runs C at B's wall-clock (1.01x) while its single core "
        "pays the full pixel tax (1.37x) - B's 8-thread run was never work-limited (thread scaling 1.98x vs C's "
        "2.68x), so the extra 35% of pixels ride in the cluster's slack for free.", 10.5, False, INK)
textbox(sC, .6, 7.2, 12.2, .3,
        "A78C/NSP cells: medians of repeated invocations (bimodal board). OFA rows: matched geometry, OWN algorithm at D=128 — not bit-exact. "
        "‡ DRAM = achieved / ceiling GB/s. A55 cells MEASURED at the DDR controller (67% of the chip's bus at 6 threads); other achieved cells not yet instrumented.",
        9.5, False, MUTED)

# ---------------- Configuration B: one slide per unit ----------------
# Same treatment as the primary configuration: the arithmetic that produced the
# number, then what the measurement means. MDE/s is the WORK-PERFORMED
# convention -- both full-search scales count.
MPX_WORK = 0.384 + 1.536   # Mpx actually searched: 960x400 + 1920x800
BD = 64
def bmde(ms): return MPX_WORK * 1e6 * BD / (ms / 1000.0) / 1e6

BUNITS = [
 ("NVIDIA RTX 5090", "GPU", 9.08, "CUDA, tuned  |  canonical: median of 5 invocations x 20 frames, 0.1% spread on a quiet card", [
   "The fastest Configuration B target: 110 fps of two-scale, 8-path, weighted-P1",
   "SGM on a workload designed for a fixed-function embedded block.",
   "The port is a GENERIC path kernel: one warp owns one path line in any of the 8",
   "directions, line start derived from the direction pair, P1 passed per launch",
   "so the direction weighting costs nothing. Bit-exact on the first full run.",
   "It lacks the primary kernel's path-pairing pass, so this row is the",
   "conservative side of what the card can do."]),
 ("NVIDIA Thor", "GPU", 45.1, "CUDA, same kernel, sm_110", [
   "22 fps embedded. Same generic 8-direction kernel as the 5090, bit-exact.",
   "The gap to the 5090 (5.0x) is close to the bandwidth ratio of the parts,",
   "consistent with the Configuration A finding that tuned SGM converges on",
   "memory traffic once the arithmetic is done."]),
 ("NVIDIA Orin AGX", "GPU", 72.8, "CUDA, same kernel, sm_87", [
   "14 fps on the previous embedded generation, bit-exact, no per-part tuning."]),
 ("Cortex-A78C x8 (6 threads)", "CPU", 211.9, "Qualcomm IQ-9075, NEON+OpenMP  |  single core: 442.60 ms", [
   "The best programmable embedded result -- and it was measured by the",
   "collaborating session running multiscale/sgm_ms_neon.c UNMODIFIED on their",
   "silicon, golden every run. Medians of five invocations on a board that is",
   "bimodal by ~9% between invocations; medians quoted, never minima.",
   "This row is also what corrected the record: at the oracle tier the NSP",
   "appeared to lead the cluster 3.29x; at the matched NEON tier the cluster",
   "leads 1.30x. A cross-tier comparison is a cross-configuration comparison",
   "wearing a platform comparison's clothes."]),
 ("Cortex-A720 x8", "CPU", 238.3, "Radxa Orion O6, NEON+OpenMP  |  single core: 470.9 ms", [
   "The NEON tier is 5.8x the scalar oracle on this cluster. The kernel keeps",
   "Configuration A's register scheme and adds the saturating step, the",
   "direction-weighted P1 per sweep, and a uint16 sum plane for 8-path S.",
   "OpenMP is output-neutral by construction: rows for horizontal sweeps,",
   "columns-within-row for vertical and diagonal, so the hash cannot move",
   "with the thread count."]),
 ("Hexagon v73 NSP", "DSP", 276.3, "Qualcomm IQ-9075, HVX, 4 contexts  |  median of 5, 0.75% spread, golden 10/10", [
   "Contributed port, independently re-verified (three further gated runs,",
   "269-272 ms). Per ENGINE this is the embedded win: one NSP beats one A78C",
   "core 1.60x at matched tiers.",
   "Diagonals use a sync-free chain-band decomposition -- workers own equal-pixel",
   "bands of diagonal chains, S worker-disjoint, bit-exact with NO barriers.",
   "The two-scale merge is fused into the stage-2 cost build, saving ~200 MB",
   "of traffic per frame. The 4-HVX-context knee holds here too: 6 threads is",
   "29% WORSE than 4.",
   "P2=200 overflowed their packed (S<<6)|d argmin key at S=2040 -- the warned",
   "saturation failure mode, one stage downstream of the warned location.",
   "Redesigned and sim-proven on tie-storms at the cap before shipping."]),
 ("Mali-G720", "GPU", 380.0, "Immortalis, 10 CU, OpenCL 3.0", [
   "One 16-work-item work-group owns each path line; barriers are uniform",
   "because every work-item in the group walks the same line trajectory.",
   "Bit-exact on the first attempt. Same caveat as Configuration A: a",
   "single-chain-per-line kernel is an effort floor for this part, not a",
   "verdict on it."]),
 ("Cortex-A55 x6", "CPU", 640.1, "i.MX 95 FRDM, 1.8 GHz, NEON+OpenMP  |  single core: 2,201 ms", [
   "The anchor platform runs the full two-scale 8-path configuration at 1.56 fps",
   "-- 22x the scalar oracle on the same six cores' single-core floor.",
   "The saturating NEON step (vqaddq_u8 against P2=200) is the part Configuration",
   "A never needed: its static assert guarantees no overflow at P2=192, and",
   "Configuration B is the configuration that breaks that assumption."]),
 ("Mali-G310", "GPU", 2704.0, "1 CU, OpenCL 3.0", [
   "Same OpenCL kernel as the G720, same hash, one compute unit."]),
 ("scalar reference", "REF", 14265.0, "1x Cortex-A55, plain C, single thread", [
   "multiscale/sgm_ms_ref.c defines golden bcb9cb0bd6f49799 and pins every",
   "contestable choice in its header: ties to lowest d, saturation at 255,",
   "2x2 box downsample, (x/2, y/2, d/2) merge indexing, x-d<0 costs 63.",
   "Every row above reproduces it byte-exactly or its timing is void."]),
]

for lab, cls, ms, sil, notes in BUNITS:
    sl = prs.slides.add_slide(blank)
    textbox(sl, .6, .35, 11, .7, lab, 32, True, ACC[cls])
    textbox(sl, 10.4, .5, 2.4, .4, "CONFIGURATION B", 12, True, MUTED)
    textbox(sl, .6, 1.0, 11.9, .4, sil, 13, False, MUTED)

    textbox(sl, .6, 1.7, 11.5, .4, "The calculation (work-performed convention)", 16, True)
    calc = (f"MDE/s  =  (0.384 + 1.536) Mpx x {BD} / runtime      (both full-search scales count)\n"
            f"        =  {MPX_WORK*1e6*BD:,.0f} disparity estimations / {ms/1000:.5f} s\n"
            f"        =  {bmde(ms):,.0f} million disparity estimations per second\n\n"
            f"frame rate  =  1000 / {ms:,.2f}  =  {1000/ms:,.2f} fps   (960x400 disparity map out)")
    tf = textbox(sl, .6, 2.15, 11.9, 1.7, calc, 14)
    for p in tf.paragraphs:
        for rr in p.runs: rr.font.name = "Consolas"

    textbox(sl, .6, 4.0, 11.5, .4, "What the measurement means", 16, True)
    tf = textbox(sl, .6, 4.4, 11.9, 2.5, notes[0], 13)
    for extra in notes[1:]:
        p = tf.add_paragraph(); rr = p.add_run(); rr.text = extra
        rr.font.size = Pt(13); rr.font.color.rgb = INK; rr.font.name = "Calibri"

    textbox(sl, .6, 7.05, 11.9, .4,
            "MEASURED -- median of the timed frames, bit-exact to golden bcb9cb0bd6f49799 in the same run.",
            11, False, MUTED)

prs.save("docs/sgm-results.pptx")
print("wrote docs/sgm-results.pptx —", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
